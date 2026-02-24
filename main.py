import base64
import io
import os
import re
import shutil
import time
import uuid
from pathlib import Path

from dotenv import load_dotenv
from fastapi import FastAPI, File, Request, UploadFile
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from google import genai
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    Image as RLImage,
    ListFlowable,
    ListItem,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)

load_dotenv()

# ── Config ────────────────────────────────────────────────────────────────────
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
gemini = genai.Client(api_key=GEMINI_API_KEY)

UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

STATIC_DIR = Path("static")
STATIC_DIR.mkdir(exist_ok=True)

# ── App ───────────────────────────────────────────────────────────────────────
app = FastAPI(title="PPTX → Teacher Guide Generator")
app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")


# ── Helpers ───────────────────────────────────────────────────────────────────

def extract_text_from_pptx(filepath: Path) -> tuple[str, str]:
    """Return (filename_stem, full_slide_text)."""
    prs = Presentation(filepath)
    lines: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
        if slide_lines:
            lines.append(f"--- Slide {slide_num} ---")
            lines.extend(slide_lines)
    return filepath.stem, "\n".join(lines)


TEACHER_GUIDE_PROMPT = """\
You are a professional curriculum designer.

Generate a complete Teacher Guide in clean semantic HTML format.

IMPORTANT RULES:
- Return ONLY HTML.
- Do NOT wrap in markdown.
- Do NOT add explanations.
- Use <h1>, <h2>, <h3>, <p>, <ul>, <li>, <strong>.
- Follow the exact hierarchy provided below.

Structure:

<h1>Session Title</h1>

<h2>Session Overview</h2>
<p>...</p>

<h2>Learning Objectives</h2>
<ul>
  <li>...</li>
</ul>

<h2>Preparation</h2>
<p>...</p>

<h2>Lesson Procedure</h2>
<h3>Initiate</h3>
<p>...</p>

<h3>Learn</h3>
<p>...</p>

<h3>Make</h3>
<p>...</p>

<h3>Share</h3>
<p>...</p>

<h2>Glossary</h2>
<ul>
  <li><strong>Term:</strong> Definition</li>
</ul>

<h2>Bonus Activities</h2>
<p>...</p>

Use clear academic language similar to an official Teacher Guide.
Audience: Write for TEACHERS, not students.

FILE NAME: {file_name}

Session content:
{slide_text}
"""


def generate_teacher_guide(file_name: str, slide_text: str) -> str:
    """Call Gemini with exponential backoff retry on 429, return HTML string."""
    prompt = TEACHER_GUIDE_PROMPT.format(
        file_name=file_name,
        slide_text=slide_text,
    )

    max_retries = 4
    delay = 10
    last_exc: Exception | None = None

    for attempt in range(max_retries):
        try:
            response = gemini.models.generate_content(
                model="gemini-2.0-flash-lite",
                contents=prompt,
            )
            raw = response.text.strip()
            # Strip accidental markdown fences
            raw = re.sub(r"^```(?:html)?\s*", "", raw, flags=re.IGNORECASE)
            raw = re.sub(r"\s*```$", "", raw)
            return raw
        except Exception as exc:
            last_exc = exc
            err_str = str(exc)
            if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
                retry_match = re.search(r"retry[\s_-]?(?:in|delay)[:\s]+([\d.]+)s", err_str, re.IGNORECASE)
                suggested = float(retry_match.group(1)) if retry_match else delay
                wait = max(suggested, delay)
                if attempt < max_retries - 1:
                    time.sleep(wait)
                    delay = wait * 2
                    continue
                raise RuntimeError(
                    "Gemini API daily quota exhausted for this API key. "
                    "Please add a new GEMINI_API_KEY to your .env file or enable billing at "
                    "https://aistudio.google.com/apikey"
                ) from exc
            raise

    raise last_exc


# ── PDF Styles ────────────────────────────────────────────────────────────────

def build_pdf_styles() -> dict:
    base = getSampleStyleSheet()
    accent = colors.HexColor("#4f46e5")

    return {
        "h1": ParagraphStyle(
            "TG_H1",
            parent=base["Normal"],
            fontSize=22,
            fontName="Helvetica-Bold",
            textColor=colors.HexColor("#1a202c"),
            spaceAfter=6,
            spaceBefore=0,
            leading=28,
        ),
        "h2": ParagraphStyle(
            "TG_H2",
            parent=base["Normal"],
            fontSize=14,
            fontName="Helvetica-Bold",
            textColor=accent,
            spaceAfter=4,
            spaceBefore=14,
            leading=18,
        ),
        "h3": ParagraphStyle(
            "TG_H3",
            parent=base["Normal"],
            fontSize=12,
            fontName="Helvetica-Bold",
            textColor=colors.HexColor("#374151"),
            spaceAfter=3,
            spaceBefore=8,
            leading=16,
        ),
        "body": ParagraphStyle(
            "TG_Body",
            parent=base["Normal"],
            fontSize=10.5,
            fontName="Helvetica",
            textColor=colors.HexColor("#374151"),
            spaceAfter=6,
            leading=16,
            alignment=TA_LEFT,
        ),
    }


def html_to_platypus(html_content: str, styles: dict) -> list:
    """Convert Quill-generated HTML to a list of ReportLab Platypus flowables."""
    flowables = []

    def clean_inner(text: str) -> str:
        """Strip Quill-specific tags, keep only <b> <i> <br/>, escape XML."""
        # Unwrap span tags (keep inner text)
        text = re.sub(r"<span\b[^>]*>(.*?)</span>", r"\1", text, flags=re.DOTALL | re.IGNORECASE)
        # Normalize strong/em to b/i
        text = re.sub(r"<strong\b[^>]*>", "<b>", text, flags=re.IGNORECASE)
        text = re.sub(r"</strong>", "</b>", text, flags=re.IGNORECASE)
        text = re.sub(r"<em\b[^>]*>", "<i>", text, flags=re.IGNORECASE)
        text = re.sub(r"</em>", "</i>", text, flags=re.IGNORECASE)
        # Normalize br
        text = re.sub(r"<br\s*/?>", "<br/>", text, flags=re.IGNORECASE)
        # Strip all remaining unknown tags, keeping inner text
        text = re.sub(r"<(?!/?(?:b|i|br/?)\b)[^>]+>", "", text, flags=re.IGNORECASE)
        # Decode common HTML entities to plain text / safe XML
        replacements = [
            ("&nbsp;", " "), ("&mdash;", "—"), ("&ndash;", "–"),
            ("&rarr;", "→"), ("&larr;", "←"), ("&ldquo;", "\u201c"),
            ("&rdquo;", "\u201d"), ("&lsquo;", "\u2018"), ("&rsquo;", "\u2019"),
            ("&hellip;", "…"), ("&copy;", "©"), ("&reg;", "®"),
        ]
        for entity, char in replacements:
            text = text.replace(entity, char)
        # Fix bare & that could break ReportLab's XML parser
        text = re.sub(r"&(?!(?:amp|lt|gt|#\d+);)", "&amp;", text)
        return text.strip()

    pos = 0
    length = len(html_content)

    while pos < length:
        ws = re.match(r"\s+", html_content[pos:])
        if ws:
            pos += ws.end()
            continue

        block = re.match(
            r"<(h1|h2|h3|p|ul|ol)\b[^>]*>(.*?)</\1>",
            html_content[pos:],
            re.DOTALL | re.IGNORECASE,
        )
        if block:
            tag = block.group(1).lower()
            inner = clean_inner(block.group(2))

            try:
                if tag == "h1":
                    flowables.append(Paragraph(inner, styles["h1"]))
                    flowables.append(Spacer(1, 0.3 * cm))

                elif tag == "h2":
                    flowables.append(Spacer(1, 0.4 * cm))
                    flowables.append(Paragraph(inner, styles["h2"]))
                    flowables.append(Spacer(1, 0.15 * cm))

                elif tag == "h3":
                    flowables.append(Spacer(1, 0.2 * cm))
                    flowables.append(Paragraph(inner, styles["h3"]))
                    flowables.append(Spacer(1, 0.1 * cm))

                elif tag == "p":
                    if inner:
                        flowables.append(Paragraph(inner, styles["body"]))
                        flowables.append(Spacer(1, 0.15 * cm))

                elif tag in ("ul", "ol"):
                    items_raw = re.findall(
                        r"<li\b[^>]*>(.*?)</li>", block.group(2), re.DOTALL | re.IGNORECASE
                    )
                    list_items = [
                        ListItem(
                            Paragraph(clean_inner(item), styles["body"]),
                            leftIndent=12,
                            spaceAfter=3,
                        )
                        for item in items_raw
                    ]
                    if list_items:
                        flowables.append(
                            ListFlowable(
                                list_items,
                                bulletType="bullet",
                                leftIndent=20,
                                bulletFontSize=8,
                            )
                        )
                        flowables.append(Spacer(1, 0.2 * cm))
            except Exception:
                pass  # skip blocks that still fail

            pos += block.end()
            continue

        # Handle <img> tags (base64 or URL)
        img_tag = re.match(
            r'<img\b[^>]*src="([^"]+)"[^>]*/?>',
            html_content[pos:],
            re.IGNORECASE,
        )
        if img_tag:
            src = img_tag.group(1)
            try:
                if src.startswith('data:image'):
                    header, b64data = src.split(',', 1)
                    img_bytes = base64.b64decode(b64data)
                    img_buf = io.BytesIO(img_bytes)
                else:
                    import urllib.request
                    with urllib.request.urlopen(src) as resp:
                        img_buf = io.BytesIO(resp.read())

                rl_img = RLImage(img_buf, width=14 * cm, kind='proportionate')
                flowables.append(Spacer(1, 0.2 * cm))
                flowables.append(rl_img)
                flowables.append(Spacer(1, 0.2 * cm))
            except Exception:
                pass
            pos += img_tag.end()
            continue

        pos += 1

    return flowables


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})


@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pptx"):
        return JSONResponse(status_code=400, content={"error": "Please upload a valid .pptx file."})

    temp_path = UPLOAD_DIR / f"{uuid.uuid4()}_{file.filename}"
    try:
        with temp_path.open("wb") as f:
            shutil.copyfileobj(file.file, f)

        file_name, slide_text = extract_text_from_pptx(temp_path)

        if not slide_text.strip():
            return JSONResponse(status_code=422, content={"error": "No readable text found in the uploaded PPTX."})

        guide_html = generate_teacher_guide(file_name, slide_text)
        return JSONResponse(content={"html": guide_html, "file_name": file_name})

    except Exception as exc:
        return JSONResponse(status_code=500, content={"error": str(exc)})
    finally:
        if temp_path.exists():
            temp_path.unlink()


@app.get("/demo")
async def demo():
    """Return a sample Teacher Guide as HTML for UI preview."""
    sample_html = """<h1>Introduction to Artificial Intelligence</h1>

<h2>Session Overview</h2>
<p>In this session, students are introduced to the concept of Artificial Intelligence through hands-on exploration and guided discussion. The lesson flows from a curiosity-driven warm-up activity into a structured explanation of how AI works, followed by a creative design task where students map out a simple AI solution to a real-world problem. By the end of the session, students will be able to articulate what AI is, how it learns, and where it appears in everyday life.</p>

<h2>Learning Objectives</h2>
<ul>
  <li>Recognise the definition of Artificial Intelligence and distinguish it from regular software.</li>
  <li>Identify at least three real-world examples of AI in everyday life.</li>
  <li>Analyse how data is used to train AI models using a simple visual example.</li>
  <li>Design a basic flowchart that describes how an AI system would solve a chosen problem.</li>
  <li>Apply critical thinking to evaluate the benefits and limitations of AI tools.</li>
</ul>

<h2>Preparation</h2>
<p><strong>Subject Knowledge:</strong> Review the basics of machine learning, supervised vs unsupervised learning, and common AI applications (image recognition, chatbots, recommendation systems).<br/><strong>You Will Need:</strong> Slides (provided), Teacher Guide, printed AI example cards, internet-connected devices for students, and access to Teachable Machine.<br/><strong>You May Need:</strong> A backup offline activity in case of connectivity issues; do a trial run of the Teachable Machine demo before class.</p>

<h2>Lesson Procedure</h2>

<h3>Initiate</h3>
<p>1. Start with a quick poll: Ask students to raise their hands if they used AI today. Accept all answers and validate them.<br/>2. Show a short 60-second video clip of AI in action (voice assistant, self-driving car, recommendation engine).<br/>3. Ask: "What do you think made the computer able to do that?" Allow 3–4 responses.<br/>4. Write the word AI on the board and say: "By the end of today, you will be able to explain exactly what this means and how it works."</p>

<h3>Learn</h3>
<p>1. Display Slide 3: Definition. Read it aloud, then ask a student to rephrase it in their own words.<br/>2. Use the analogy: "Teaching an AI is like teaching a child — you show it many examples until it learns the pattern."<br/>3. Slide 4 — Show the training data diagram. Walk through each step: Data → Training → Model → Prediction.<br/>4. Pause and ask: "If we trained an AI on only pictures of cats, what would happen if we showed it a dog?" Wait for responses.<br/>5. Slide 5 — Real-world examples. For each example, ask students where they have seen it.</p>

<h3>Make</h3>
<p>1. Distribute the "AI Problem Solver" worksheet (or project it on screen).<br/>2. Students choose a real problem (e.g., detecting ripe fruit, filtering spam emails).<br/>3. They draw a simple 4-step flowchart: Problem → Data Needed → How AI Learns → Output.<br/>4. Circulate the room. For students who are stuck, ask: "What information would a human need to solve this? That is likely the data your AI needs."<br/>5. Early finishers: Challenge them to add a "What could go wrong?" box to their flowchart.</p>

<h3>Share</h3>
<p>1. Ask 2–3 volunteers to share their flowcharts with the class.<br/>2. For each presentation, prompt the audience: "Does this AI need a lot of data or a little? Why?"<br/>3. Wrap up by asking the full class: "Name one thing AI is great at and one thing humans still do better."<br/>4. Close with: "Next lesson we will actually train our own AI model using Teachable Machine."</p>

<h2>Glossary</h2>
<ul>
  <li><strong>Artificial Intelligence (AI):</strong> A computer system that can perform tasks that normally require human intelligence, like recognising images or understanding speech.</li>
  <li><strong>Machine Learning:</strong> A type of AI where the computer learns from examples (data) instead of being told exact rules.</li>
  <li><strong>Training Data:</strong> A large collection of examples used to teach an AI model how to make decisions.</li>
  <li><strong>Algorithm:</strong> A step-by-step set of instructions a computer follows to solve a problem.</li>
  <li><strong>Model:</strong> The result of training an AI — a program that can make predictions based on what it has learned.</li>
</ul>

<h2>Bonus Activities</h2>
<p>1. Explore Google's Teachable Machine (teachablemachine.withgoogle.com) and train a model to recognise hand gestures using only the webcam.<br/>2. Research one AI system that had a bias problem (e.g., facial recognition accuracy disparities) and write a short paragraph explaining what went wrong and how it could be fixed.<br/>3. Create a short comic strip (4 panels) telling the story of how an AI learns to do one task.</p>"""

    return JSONResponse(content={"html": sample_html, "file_name": "Sample_Lesson_Introduction_to_AI"})


@app.post("/export-pdf")
async def export_pdf(request: Request):
    """Accept HTML content and return a professionally styled PDF."""
    body = await request.json()
    html_content = body.get("html", "")
    file_name = body.get("file_name", "teacher_guide")

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=2.5 * cm,
        leftMargin=2.5 * cm,
        topMargin=2.5 * cm,
        bottomMargin=2.5 * cm,
        title=file_name,
        author="Teacher Guide Generator",
    )

    styles = build_pdf_styles()
    flowables = html_to_platypus(html_content, styles)

    if not flowables:
        flowables = [Paragraph("No content to export.", styles["body"])]

    doc.build(flowables)
    buffer.seek(0)

    safe_name = re.sub(r"[^\w\-]", "_", file_name)
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}.pdf"'},
    )
